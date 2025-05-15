from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
import textwrap

def prevent_text_overflow(text_frame_or_paragraph, text_content, max_font_size=32, min_font_size=10, width=None):
    """
    Prevents text overflow by automatically adjusting font size and adding line breaks.
    
    Args:
        text_frame_or_paragraph: The text frame or paragraph object from python-pptx
        text_content (str): The text to be added
        max_font_size (int): Maximum font size to start with
        min_font_size (int): Minimum allowed font size
        width (float, optional): Width constraint in inches. If None, uses text frame width
    
    Returns:
        bool: True if text was successfully fitted, False if it couldn't be fitted
    """
    if not text_content:
        return True

    # Determine if we're working with a text frame or paragraph
    is_paragraph = not hasattr(text_frame_or_paragraph, 'paragraphs')
    
    if is_paragraph:
        paragraph = text_frame_or_paragraph
    else:
        text_frame_or_paragraph.clear()
        paragraph = text_frame_or_paragraph.paragraphs[0]
    
    # Set initial formatting
    paragraph.alignment = PP_ALIGN.LEFT
    font = paragraph.font
    font.name = 'Calibri'
    font.color.rgb = RGBColor(0, 0, 0)
    
    # If width is not provided, try to get it from the text frame
    if width is None:
        try:
            if is_paragraph:
                width = Inches(8)  # Default width for paragraphs
            else:
                width = text_frame_or_paragraph.width
        except AttributeError:
            width = Inches(8)  # Default slide width
    
    # Calculate approximate characters per line (assuming average char width)
    avg_char_width_pts = max_font_size * 0.6  # Approximate character width
    max_chars_per_line = int((width / Pt(1)) / avg_char_width_pts)
    
    # Try different font sizes from max to min
    current_font_size = max_font_size
    while current_font_size >= min_font_size:
        font.size = Pt(current_font_size)
        
        # Adjust chars per line based on current font size
        chars_per_line = int(max_chars_per_line * (max_font_size / current_font_size))
        
        # Wrap text
        wrapped_lines = textwrap.wrap(text_content, width=chars_per_line)
        
        # For bullet points, we'll be more conservative with vertical space
        if len(wrapped_lines) <= 3:  # Allow up to 3 lines per bullet point
            paragraph.text = '\n'.join(wrapped_lines)
            return True
        
        current_font_size -= 2
    
    # If we get here, text couldn't be fitted even at minimum font size
    # Truncate the text and add ellipsis
    font.size = Pt(min_font_size)
    chars_per_line = int(max_chars_per_line * (max_font_size / min_font_size))
    wrapped_lines = textwrap.wrap(text_content, width=chars_per_line)[:2]  # Take first 2 lines
    
    if wrapped_lines:
        last_line = wrapped_lines[-1]
        if len(last_line) > chars_per_line - 3:
            wrapped_lines[-1] = last_line[:chars_per_line-3] + '...'
        else:
            wrapped_lines[-1] = last_line + '...'
    
    paragraph.text = '\n'.join(wrapped_lines)
    return False

def format_bullet_points(text_frame, bullet_points, level_font_sizes=None):
    """
    Formats a list of bullet points with proper indentation and font sizes.
    
    Args:
        text_frame: The text frame object from python-pptx
        bullet_points (list): List of strings or tuples (text, level)
        level_font_sizes (dict, optional): Dictionary mapping level to font size
    """
    if level_font_sizes is None:
        level_font_sizes = {
            0: 24,  # Main points
            1: 20,  # Sub points
            2: 18   # Sub-sub points
        }
    
    # Enable auto-size for the text frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.word_wrap = True
    text_frame.clear()
    
    for point in bullet_points:
        if isinstance(point, tuple):
            text, level = point
        else:
            text, level = point, 0
            
        # Add paragraph and set its properties
        paragraph = text_frame.add_paragraph()
        paragraph.level = level
        paragraph.text = text  # Set text first
        
        # Set font properties
        font = paragraph.font
        font.size = Pt(level_font_sizes.get(level, 18))
        font.name = 'Calibri'
        font.color.rgb = RGBColor(0, 0, 0)
        
        # Enable bullets
        try:
            paragraph._pPr.get_or_add_pPr().set('bullet', '•')
        except AttributeError:  # If _pPr is not available
            try:
                paragraph._element.get_or_add_pPr().set('bullet', '•')
            except (AttributeError, KeyError):  # If element access or bullet setting fails
                # Skip bullet setting if both methods fail
                pass
        
        # Handle text overflow by wrapping and truncating if necessary
        if len(text) > 100:  # Only process long text
            prevent_text_overflow(paragraph, text,
                                max_font_size=level_font_sizes.get(level, 18),
                                min_font_size=10) 