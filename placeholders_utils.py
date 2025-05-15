from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

def get_placeholder_details(slide_layout):
    """
    Extracts details (ID and type name) for all placeholders in a given slide layout.
    """
    placeholder_details = []
    placeholders = slide_layout.placeholders
    if placeholders:
        for placeholder in placeholders:
            # Access placeholder type using placeholder_format.type
            placeholder_type = placeholder.placeholder_format.type
            # Get the enum member name for readability
            type_name = PP_PLACEHOLDER(placeholder_type).name if placeholder_type in PP_PLACEHOLDER.__members__.values() else "UNKNOWN"
            placeholder_details.append({
                'id': placeholder.placeholder_format.idx,
                'type_name': type_name
            })
    return placeholder_details

def _find_id_by_types(placeholder_details, type_names_in_order):
    """
    Helper function to find the ID of the first placeholder matching a list
    of type names, searched in the provided order. Case-insensitive.
    """
    for type_name in type_names_in_order:
        for detail in placeholder_details:
            # Case-insensitive comparison
            if detail['type_name'].upper() == type_name.upper():
                return detail['id']
    return None

def _find_indexed_id_by_types(placeholder_details, type_names_in_order, index):
    """
    Helper function to find the ID of the placeholder matching a list of types
    at a specific index (0-based) among matches, sorted by placeholder ID (idx).
    """
    all_matches_sorted = []
    target_types_upper = [t.upper() for t in type_names_in_order]

    # Collect all placeholders matching the target types
    for detail in placeholder_details:
        if detail['type_name'].upper() in target_types_upper:
             all_matches_sorted.append({'id': detail['id'], 'type_name': detail['type_name']})

    # Sort by placeholder index (id)
    all_matches_sorted.sort(key=lambda x: x['id'])

    if index < len(all_matches_sorted):
        return all_matches_sorted[index]['id']
    # Fallback for index 1 if only one match exists (e.g., requesting right_content when only one content box found)
    elif index > 0 and len(all_matches_sorted) == 1:
         print(f"Warning: Requested index {index} for types {type_names_in_order}, but only one match found. Falling back to index 0.")
         return all_matches_sorted[0]['id']
    else:
        # Return None if index is out of bounds and not the specific fallback case above
        return None


def find_placeholder_id( placeholder_details,content_key):
    """
    Finds the most probable placeholder ID for a given content key based on
    placeholder types and common layout structures.

    Args:
        content_key (str): The key representing the content type (e.g., "title", "content").
        placeholder_details (list): A list of dicts with 'id' and 'type_name' for placeholders.

    Returns:
        int or None: The ID (idx) of the most probable placeholder, or None if not found.
    """
    if not placeholder_details:
        return None

    key = content_key.lower() # Normalize key

    # --- Title Placeholders ---
    if key in ["title", "section_title"]:
        # Prefer TITLE, CENTER_TITLE. Fallback to first BODY if necessary.
        return _find_id_by_types(placeholder_details, ["TITLE", "CENTER_TITLE", "BODY"])

    # --- Subtitle Placeholders ---
    elif key in ["subtitle", "section_description"]:
        # Prefer SUBTITLE. Fallback to first BODY.
        return _find_id_by_types(placeholder_details, ["SUBTITLE", "BODY"])

    # --- Main Content Placeholders (Single) ---
    elif key == "content":
        # Prefer BODY, CONTENT, OBJECT. Find the first one.
        # Use indexed search with index 0 to ensure we get the first one based on ID sorting.
        return _find_indexed_id_by_types(placeholder_details, ["BODY", "CONTENT", "OBJECT"], 0)

    # --- Two Content / Comparison Placeholders ---
    elif key in ["left_content", "left_comparison_content", "left_heading"]:
        # Find the first BODY, CONTENT, or OBJECT placeholder (sorted by ID).
        # Assuming 'left' corresponds to the lower placeholder ID.
        return _find_indexed_id_by_types(placeholder_details, ["BODY", "CONTENT", "OBJECT"], 0)
    elif key in ["right_content", "right_comparison_content", "right_heading"]:
        # Find the second BODY, CONTENT, or OBJECT placeholder (sorted by ID).
        # Assuming 'right' corresponds to the higher placeholder ID.
        # Includes fallback to the first if only one exists.
        return _find_indexed_id_by_types(placeholder_details, ["BODY", "CONTENT", "OBJECT"], 1)

    # --- Caption/Description Placeholders ---
    elif key == "caption_text":
        # Often a BODY or OBJECT placeholder. Let's try finding BODY/CONTENT/OBJECT.
        # This might conflict if there's also a main 'content'. Needs layout context.
        # Simple approach: Find first BODY/CONTENT/OBJECT (same as 'content').
        return _find_indexed_id_by_types(placeholder_details, ["BODY", "CONTENT", "OBJECT"], 0)
    elif key == "object_description":
        # Prefer OBJECT, PICTURE, then BODY.
        return _find_id_by_types(placeholder_details, ["OBJECT", "PICTURE", "BODY"])
    elif key == "picture_description":
        # Prefer PICTURE, OBJECT, then BODY.
        return _find_id_by_types(placeholder_details, ["PICTURE", "OBJECT", "BODY"])
    else:
        print(f"Warning: Unknown content key '{content_key}'. Cannot determine placeholder.")
        return None

def print_layout_details(pptx_path):
    """
    Loads a PowerPoint presentation and prints details about its slide layouts
    and their placeholders.
    """
    try:
        prs = Presentation(pptx_path)
        print(f"Loaded presentation: {pptx_path}")
        print("-" * 30)

        if not prs.slide_layouts:
            print("No slide layouts found in this presentation.")
            return

        print(f"Found {len(prs.slide_layouts)} slide layouts:")

        for i, slide_layout in enumerate(prs.slide_layouts):
            print(f"\nLayout {i + 1}: {slide_layout.name}")
            details = get_placeholder_details(slide_layout)
            print(f"  Placeholders: {len(details)}")
            if details:
                for idx, detail in enumerate(details):
                    print(f"    - Placeholder {idx + 1}: Type = {detail['type_name']} (ID: {detail['id']})")
            else:
                print("    - No placeholders in this layout.")
            print("-" * 20)

    except Exception as e:
        print(f"An error occurred: {e}")

# if __name__ == "__main__":
#     # Replace with the actual path to your PowerPoint file
#     presentation_path = "./templates/Organic presentation.pptx"
#     # print_layout_details(presentation_path) # Original function call (optional)

#     # --- Test the new find_placeholder_id function ---
#     try:
#         prs = Presentation(presentation_path)
#         print(f"\n--- Testing find_placeholder_id ---")

#         if not prs.slide_layouts:
#             print("No slide layouts found.")
#         else:
#             # Test specific layouts by index
#             layouts_to_test = {
#                 0: ["title", "subtitle"], # Title Slide
#                 1: ["title", "content"], # Title and Content
#                 3: ["title", "left_content", "right_content"], # Two Content
#                 # Add more layouts and keys to test as needed
#                 # 4: ["title", "left_heading", "right_heading", "left_comparison_content", "right_comparison_content"], # Comparison
#                 # 7: ["title", "caption_text", "object_description"], # Content with Caption
#             }

#             for layout_idx, keys_to_find in layouts_to_test.items():
#                  if layout_idx < len(prs.slide_layouts):
#                      test_layout = prs.slide_layouts[layout_idx]
#                      print(f"\nTesting Layout {layout_idx}: {test_layout.name}")
#                      details = get_placeholder_details(test_layout)
#                      print(f"  Placeholders found: {details}")

#                      if not details:
#                          print("  No placeholders in this layout to test.")
#                          continue

#                      for key in keys_to_find:
#                          found_id = find_placeholder_id(key, details)
#                          print(f"  Probable ID for '{key}': {found_id}")
#                  else:
#                      print(f"\nLayout index {layout_idx} out of range (max: {len(prs.slide_layouts)-1}).")


#     except Exception as e:
#         print(f"An error occurred during testing: {e}")
