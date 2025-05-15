import json
import os
import re
import traceback
from typing import Any, Dict
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

# --- Default Configuration ---
DEFAULT_CONFIG = {
    # Core
    'default_font_name': "Calibri",
    'default_font_color_rgb': "000000",
    # Font Sizes
    'default_title_font_size_pt': 36,
    'default_body_font_size_pt': 20,
    'default_subtitle_font_size_pt': 24,
    'default_caption_font_size_pt': 14,
    'default_notes_font_size_pt': 10,
    # Background
    'master_background_color_rgb': "", # e.g., "FFFFFF" for white, empty means use template default
    # Alignment (Specific keys for each placeholder in each layout)
    'alignment_layout_0_title': "CENTER",
    'alignment_layout_0_subtitle': "CENTER",
    'alignment_layout_1_title': "LEFT",
    'alignment_layout_1_body': "LEFT",
    'alignment_layout_2_title': "LEFT",
    'alignment_layout_2_body': "LEFT",
    'alignment_layout_3_title': "LEFT",
    'alignment_layout_3_left_body': "LEFT",
    'alignment_layout_3_right_body': "LEFT",
    'alignment_layout_4_title': "LEFT",
    'alignment_layout_4_left_head': "LEFT",
    'alignment_layout_4_right_head': "LEFT",
    'alignment_layout_4_left_body': "LEFT",
    'alignment_layout_4_right_body': "LEFT",
    'alignment_layout_5_title': "CENTER",
    'alignment_layout_7_title': "LEFT",
    'alignment_layout_7_caption': "LEFT",
    'alignment_layout_7_body': "LEFT",
    'alignment_layout_8_title': "LEFT",
    'alignment_layout_8_caption': "LEFT",
    'alignment_notes': "LEFT",
    # Vertical Anchor
    'default_vertical_anchor': "TOP", # Options: TOP, MIDDLE, BOTTOM
    # Dynamic Sizing & Formatting
    'enable_dynamic_body_font_size': True, # Reduce body font size if content is long
    'smaller_content_font_size_pt': 16, # Font size to use when dynamic sizing is triggered
    'dynamic_size_item_count_threshold': 6, # Trigger smaller size if more list items than this
    'dynamic_size_char_count_threshold': 400, # Trigger smaller size if total characters exceed this
    'indent_level_font_size_reduction_pt': 2, # Reduce font size by this amount per indent level
    'min_font_size_pt': 10, # Minimum font size after reductions
}

# --- Helper Functions ---

def parse_rgb(rgb_string):
    """Converts a 6-digit hex string (e.g., "FF0000") to an RGBColor object."""
    if rgb_string and isinstance(rgb_string, str) and len(rgb_string) == 6:
        try:
            r = int(rgb_string[0:2], 16)
            g = int(rgb_string[2:4], 16)
            b = int(rgb_string[4:6], 16)
            return RGBColor(r, g, b)
        except ValueError:
            print(f"Warning: Invalid RGB hex string format: '{rgb_string}'.")
            return None
    elif rgb_string: # Catch non-string or incorrect length
        print(f"Warning: Invalid RGB value: '{rgb_string}'. Expected 6-digit hex string.")
    return None

def get_alignment_enum(align_str):
    """Maps alignment string (e.g., "CENTER") to a PP_ALIGN enum."""
    mapping = {
        "LEFT": PP_ALIGN.LEFT,
        "CENTER": PP_ALIGN.CENTER,
        "RIGHT": PP_ALIGN.RIGHT,
        "JUSTIFY": PP_ALIGN.JUSTIFY,
        "DISTRIBUTE": PP_ALIGN.DISTRIBUTE,
        "THAI_DISTRIBUTE": PP_ALIGN.THAI_DISTRIBUTE
    }
    # Default to LEFT if invalid or not specified
    return mapping.get(str(align_str).upper(), PP_ALIGN.LEFT)

def get_vertical_anchor_enum(anchor_str):
    """Maps vertical anchor string (e.g., "MIDDLE") to an MSO_ANCHOR enum."""
    mapping = {
        "TOP": MSO_ANCHOR.TOP,
        "MIDDLE": MSO_ANCHOR.MIDDLE,
        "BOTTOM": MSO_ANCHOR.BOTTOM,
    }
    # Default to TOP if invalid or not specified
    return mapping.get(str(anchor_str).upper(), MSO_ANCHOR.TOP)

# --- Rich Text Formatting Helpers ---

def add_formatted_run(paragraph, text, bold=False, italic=False, underline=False):
    """Adds a text run with specified formatting to a paragraph."""
    if not text:
        return None # Avoid adding empty runs
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.bold = bold
    font.italic = italic
    font.underline = underline
    # Base color, name, and size are set by the calling function
    return run

def parse_and_add_formatted_text(paragraph, text_string, base_font_size_pt, base_font_name, base_font_color_rgb, config):
    """
    Parses simple markdown (**bold**, *italic*, <u>underline</u>) in a string
    and adds formatted runs to the paragraph. Applies base font properties and
    adjusts size based on indent level.
    """
    # Regex to split text by markdown markers, keeping the markers
    pattern = r'(\*\*.*?\*\*|\*.*?\*|<u>.*?</u>)'
    parts = re.split(pattern, text_string)

    # Adjust font size based on paragraph level (indentation)
    indent_reduction = config.get('indent_level_font_size_reduction_pt', 2)
    min_font_size = config.get('min_font_size_pt', 10)
    effective_font_size_pt = max(
        base_font_size_pt - (paragraph.level * indent_reduction),
        min_font_size
    )

    base_color = parse_rgb(base_font_color_rgb)

    for part in parts:
        if not part:
            continue

        # Determine formatting based on start/end markers
        is_bold = part.startswith('**') and part.endswith('**') and len(part) > 4
        is_italic = part.startswith('*') and part.endswith('*') and len(part) > 2
        is_underline = part.startswith('<u>') and part.endswith('</u>') and len(part) > 7

        run_text = part
        run_bold = False
        run_italic = False
        run_underline = False

        # Extract text content and set format flags
        if is_bold:
            run_text = part[2:-2]
            run_bold = True
        elif is_italic:
            run_text = part[1:-1]
            run_italic = True
        elif is_underline:
            run_text = part[3:-4]
            run_underline = True

        # Add the run with detected formatting
        run = add_formatted_run(paragraph, run_text, bold=run_bold, italic=run_italic, underline=run_underline)

        # Apply base font properties (name, size, color) to the new run
        if run:
            run.font.size = Pt(effective_font_size_pt)
            run.font.name = base_font_name
            if base_color:
                run.font.color.rgb = base_color

# --- Text Frame Population ---
def populate_text_frame(text_frame, content_data, base_font_size_pt, base_font_name, base_font_color_rgb, config,
                        alignment=None, vertical_anchor=None):
    """
    Populates a text frame with content (string or list), applying formatting,
    alignment, vertical anchor, and dynamic font sizing based on config.
    """
    text_frame.clear()
    text_frame.word_wrap = True

    final_vertical_anchor_str = vertical_anchor if vertical_anchor else config.get('default_vertical_anchor', 'TOP')
    text_frame.vertical_anchor = get_vertical_anchor_enum(final_vertical_anchor_str)

    if not content_data:
        return

    content_list = content_data if isinstance(content_data, list) else [content_data]

    effective_base_font_size = base_font_size_pt
    if config.get('enable_dynamic_body_font_size', True) and base_font_size_pt == config['default_body_font_size_pt']:
        non_empty_content = [item for item in content_list if isinstance(item, str) and item.strip()]
        num_content_items = len(non_empty_content)
        total_content_length = sum(len(item.strip()) for item in non_empty_content)

        item_threshold = config.get('dynamic_size_item_count_threshold', 6)
        char_threshold = config.get('dynamic_size_char_count_threshold', 400)
        smaller_size = config.get('smaller_content_font_size_pt', 16)

        if num_content_items > item_threshold or total_content_length > char_threshold:
            effective_base_font_size = smaller_size

    final_alignment_str = alignment if alignment else config.get('alignment_layout_1_body', 'LEFT')
    final_alignment_enum = get_alignment_enum(final_alignment_str)

    for item_text in content_list:
        if not isinstance(item_text, str):
            print(f"      Warning: Skipping non-string item in content list: {type(item_text)}")
            continue

        item_text_stripped = item_text.lstrip()
        if not item_text_stripped and len(item_text) > 0:
             continue
        elif not item_text_stripped:
             continue

        p = text_frame.add_paragraph()

        leading_spaces = len(item_text) - len(item_text_stripped)
        p.level = min(leading_spaces // 2, 5)

        p.alignment = final_alignment_enum

        parse_and_add_formatted_text(p, item_text_stripped,
                                     effective_base_font_size,
                                     base_font_name,
                                     base_font_color_rgb,
                                     config)

# --- Placeholder Finding Logic ---
def find_placeholder(slide, placeholder_type, hint_name=None, hint_idx=None):
    """
    Finds a placeholder on a slide using its type and optional hints (name substring, index).
    placeholder_type: e.g., PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.BODY
    hint_name: A string expected in the placeholder's internal name.
    hint_idx: The expected placeholder index (e.g., 0 for title).
    """
    # Try finding by specific PP_PLACEHOLDER type (most reliable)
    if isinstance(placeholder_type, PP_PLACEHOLDER):
        for shape in slide.placeholders:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == placeholder_type:
                return shape

    # Try finding by general MSO_SHAPE_TYPE and hints if PP_PLACEHOLDER type didn't work
    found_by_type = []
    for shape in slide.placeholders:
        if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
             if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == placeholder_type:
                 found_by_type.append(shape)

    if found_by_type:
        # If multiple match the type, use hints to disambiguate
        if hint_name:
            for shape in found_by_type:
                if shape.name and hint_name in shape.name:
                    return shape
        if hint_idx is not None:
             for shape in found_by_type:
                 if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx == hint_idx:
                    return shape
        # If hints didn't help, return the first one found matching the type
        return found_by_type[0]

    # Fallback: Try finding by hint_idx only
    if hint_idx is not None:
        for shape in slide.placeholders:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx == hint_idx:
                return shape

    # Fallback: Try finding by hint_name only
    if hint_name:
        for shape in slide.placeholders:
            if shape.name and hint_name in shape.name:
                return shape

    return None


# --- Presentation Creation Function ---
def create_ppt(slides_data, output_filename, config):
    """Creates a PowerPoint presentation from a list of slide data dictionaries."""
    if not isinstance(slides_data, list):
        print(f"Error: Input 'slides_data' must be a list, but received {type(slides_data)}. Check JSON structure.")
        return
    if not slides_data:
        print("Error: No slide data provided (list is empty).")
        return

    prs = Presentation()
    print("Created presentation using python-pptx default template.")

    if config.get('master_background_color_rgb'):
        bg_color = parse_rgb(config['master_background_color_rgb'])
        if bg_color:
            try:
                # Apply to the first slide master (usually the main one)
                if prs.slide_masters:
                    fill = prs.slide_masters[0].background.fill
                    fill.solid()
                    fill.fore_color.rgb = bg_color
                    print(f"Set master background color to {config['master_background_color_rgb']}")
                else:
                    print("Warning: Could not set master background color - no slide masters found.")
            except Exception as e:
                print(f"Warning: Failed to set master background color: {e}")

    # Map JSON 'layout_idx' to the default python-pptx template layout indices
    # This mapping assumes the standard default template. Custom templates might differ.
    LAYOUT_MAP = {
        0: 0, # Title Slide
        1: 1, # Title and Content
        2: 2, # Section Header
        3: 3, # Two Content
        4: 4, # Comparison
        5: 5, # Title Only
        6: 6, # Blank
        7: 7, # Content with Caption
        8: 8, # Picture with Caption
    }
    layouts = prs.slide_layouts

    if not layouts:
        print("Error: No slide layouts found in the presentation template. Cannot proceed.")
        return

    # --- Process Each Slide ---
    for i, slide_info in enumerate(slides_data):
        print(f"Processing slide {i+1}...")
        if not isinstance(slide_info, dict):
            print(f"  Warning: Slide {i+1} data is not a dictionary ({type(slide_info)}). Skipping.")
            continue

        try:
            layout_idx_from_json = slide_info.get('layout_idx')
            if layout_idx_from_json is None:
                print(f"  Error: Slide {i+1} is missing 'layout_idx'. Skipping.")
                continue

            try:
                layout_idx = int(layout_idx_from_json)
                if layout_idx not in LAYOUT_MAP:
                    print(f"  Error: Slide {i+1} has unsupported 'layout_idx': {layout_idx}. Skipping.")
                    continue
            except (ValueError, TypeError):
                print(f"  Error: Slide {i+1} has non-integer 'layout_idx': '{layout_idx_from_json}'. Skipping.")
                continue

            target_layout_index = LAYOUT_MAP[layout_idx]
            if not (0 <= target_layout_index < len(layouts)):
                 print(f"  Error: Mapped layout index {target_layout_index} (from JSON index {layout_idx}) is out of range for available layouts ({len(layouts)}). Skipping slide.")
                 continue

            slide_layout = layouts[target_layout_index]
            print(f"  Using layout index: {layout_idx} (template layout {target_layout_index}: '{slide_layout.name}')")

            slide = prs.slides.add_slide(slide_layout)

            default_vert_anchor = config.get('default_vertical_anchor', 'TOP')

            # --- Populate Placeholders Based on Layout ---
            title_text = slide_info.get('title', '')
            notes_text = slide_info.get('notes', '')

            # Layout-specific population using find_placeholder and populate_text_frame
            if layout_idx == 0: # Title Slide
                title_ph = find_placeholder(slide, PP_PLACEHOLDER.TITLE, hint_idx=0)
                subtitle_ph = find_placeholder(slide, PP_PLACEHOLDER.SUBTITLE, hint_idx=1)
                if title_ph and title_text:
                    align = config.get('alignment_layout_0_title', 'CENTER')
                    populate_text_frame(title_ph.text_frame, title_text,
                                        config['default_title_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config,
                                        alignment=align, vertical_anchor=default_vert_anchor)
                subtitle_text = slide_info.get('subtitle')
                if subtitle_ph and subtitle_text:
                    align = config.get('alignment_layout_0_subtitle', 'CENTER')
                    populate_text_frame(subtitle_ph.text_frame, subtitle_text,
                                        config['default_subtitle_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config,
                                        alignment=align, vertical_anchor=default_vert_anchor)

            elif layout_idx == 1: # Title and Content
                title_ph = find_placeholder(slide, PP_PLACEHOLDER.TITLE, hint_idx=0)
                content_ph = find_placeholder(slide, PP_PLACEHOLDER.BODY, hint_idx=1)
                if title_ph and title_text:
                    align = config.get('alignment_layout_1_title', 'LEFT')
                    populate_text_frame(title_ph.text_frame, title_text,
                                        config['default_title_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config,
                                        alignment=align, vertical_anchor=default_vert_anchor)
                content_data = slide_info.get('content')
                if content_ph and content_data:
                    align = config.get('alignment_layout_1_body', 'LEFT')
                    populate_text_frame(content_ph.text_frame, content_data,
                                        config['default_body_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config,
                                        alignment=align, vertical_anchor=default_vert_anchor)

            elif layout_idx == 2: # Section Header
                # Section headers often use Title and Body placeholders differently
                section_title_ph = find_placeholder(slide, PP_PLACEHOLDER.TITLE, hint_idx=0) # Often the main text
                section_desc_ph = find_placeholder(slide, PP_PLACEHOLDER.BODY, hint_idx=1) # Often smaller subtitle text

                section_title = slide_info.get('section_title')
                if section_title_ph and section_title:
                    align = config.get('alignment_layout_2_title', 'LEFT')
                    populate_text_frame(section_title_ph.text_frame, section_title, config['default_title_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config, alignment=align, vertical_anchor=default_vert_anchor)
                section_desc = slide_info.get('section_description')
                if section_desc_ph and section_desc:
                    align = config.get('alignment_layout_2_body', 'LEFT')
                    populate_text_frame(section_desc_ph.text_frame, section_desc, config['default_body_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config, alignment=align, vertical_anchor=default_vert_anchor)

            elif layout_idx == 3: # Two Content
                title_ph = find_placeholder(slide, PP_PLACEHOLDER.TITLE, hint_idx=0)
                left_ph = find_placeholder(slide, PP_PLACEHOLDER.BODY, hint_idx=1)
                right_ph = find_placeholder(slide, PP_PLACEHOLDER.BODY, hint_idx=2)

                if title_ph and title_text:
                    align = config.get('alignment_layout_3_title', 'LEFT')
                    populate_text_frame(title_ph.text_frame, title_text, config['default_title_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config, alignment=align, vertical_anchor=default_vert_anchor)
                left_content = slide_info.get('left_content')
                if left_ph and left_content:
                    align = config.get('alignment_layout_3_left_body', 'LEFT')
                    populate_text_frame(left_ph.text_frame, left_content, config['default_body_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config, alignment=align, vertical_anchor=default_vert_anchor)
                right_content = slide_info.get('right_content')
                if right_ph and right_content:
                    align = config.get('alignment_layout_3_right_body', 'LEFT')
                    populate_text_frame(right_ph.text_frame, right_content, config['default_body_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config, alignment=align, vertical_anchor=default_vert_anchor)

            elif layout_idx == 4: # Comparison
                title_ph = find_placeholder(slide, PP_PLACEHOLDER.TITLE, hint_idx=0)
                # Use hints as comparison layout indices can vary
                left_head_ph = find_placeholder(slide, PP_PLACEHOLDER.BODY, hint_idx=1, hint_name="Head")
                right_head_ph = find_placeholder(slide, PP_PLACEHOLDER.BODY, hint_idx=2, hint_name="Head")
                left_body_ph = find_placeholder(slide, PP_PLACEHOLDER.BODY, hint_idx=3, hint_name="Text")
                right_body_ph = find_placeholder(slide, PP_PLACEHOLDER.BODY, hint_idx=4, hint_name="Text")

                if title_ph and title_text:
                    align = config.get('alignment_layout_4_title', 'LEFT')
                    populate_text_frame(title_ph.text_frame, title_text, config['default_title_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config, alignment=align, vertical_anchor=default_vert_anchor)
                left_heading = slide_info.get('left_heading')
                if left_head_ph and left_heading:
                    align = config.get('alignment_layout_4_left_head', 'LEFT')
                    populate_text_frame(left_head_ph.text_frame, left_heading, config['default_body_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config, alignment=align, vertical_anchor=default_vert_anchor)
                right_heading = slide_info.get('right_heading')
                if right_head_ph and right_heading:
                    align = config.get('alignment_layout_4_right_head', 'LEFT')
                    populate_text_frame(right_head_ph.text_frame, right_heading, config['default_body_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config, alignment=align, vertical_anchor=default_vert_anchor)
                left_comp_content = slide_info.get('left_comparison_content')
                if left_body_ph and left_comp_content:
                    align = config.get('alignment_layout_4_left_body', 'LEFT')
                    populate_text_frame(left_body_ph.text_frame, left_comp_content, config['default_body_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config, alignment=align, vertical_anchor=default_vert_anchor)
                right_comp_content = slide_info.get('right_comparison_content')
                if right_body_ph and right_comp_content:
                    align = config.get('alignment_layout_4_right_body', 'LEFT')
                    populate_text_frame(right_body_ph.text_frame, right_comp_content, config['default_body_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config, alignment=align, vertical_anchor=default_vert_anchor)

            elif layout_idx == 5: # Title Only
                title_ph = find_placeholder(slide, PP_PLACEHOLDER.TITLE, hint_idx=0)
                if title_ph and title_text:
                    align = config.get('alignment_layout_5_title', 'CENTER')
                    # Override vertical anchor to MIDDLE for better centering
                    vert_anchor = MSO_ANCHOR.MIDDLE
                    populate_text_frame(title_ph.text_frame, title_text,
                                        config['default_title_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config,
                                        alignment=align, vertical_anchor=vert_anchor)

            elif layout_idx == 6: # Blank
                print("  Blank layout selected. No standard content placeholders.")
                if title_text:
                    print(f"  Warning (Layout {layout_idx}): Title text provided for a blank layout. It will not be added automatically.")

            elif layout_idx == 7: # Content with Caption
                title_ph = find_placeholder(slide, PP_PLACEHOLDER.TITLE, hint_idx=0)
                content_ph = find_placeholder(slide, PP_PLACEHOLDER.BODY, hint_idx=1)
                caption_ph = find_placeholder(slide, PP_PLACEHOLDER.BODY, hint_idx=2, hint_name="Caption")

                if title_ph and title_text:
                    align = config.get('alignment_layout_7_title', 'LEFT')
                    populate_text_frame(title_ph.text_frame, title_text, config['default_title_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config, alignment=align, vertical_anchor=default_vert_anchor)
                caption_text = slide_info.get('caption_text')
                if caption_ph and caption_text:
                    align = config.get('alignment_layout_7_caption', 'LEFT')
                    populate_text_frame(caption_ph.text_frame, caption_text, config['default_caption_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config, alignment=align, vertical_anchor=default_vert_anchor) # Use caption size
                obj_desc = slide_info.get('object_description') # Text for the main content area
                if content_ph and obj_desc:
                    # Populate content placeholder if it's empty
                    if not content_ph.has_text_frame or not content_ph.text_frame.text:
                         align = config.get('alignment_layout_7_body', 'LEFT')
                         populate_text_frame(content_ph.text_frame, obj_desc, config['default_body_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config, alignment=align, vertical_anchor=default_vert_anchor)

            elif layout_idx == 8: # Picture with Caption
                title_ph = find_placeholder(slide, PP_PLACEHOLDER.TITLE, hint_idx=0)
                picture_ph = find_placeholder(slide, PP_PLACEHOLDER.PICTURE, hint_idx=1)
                caption_ph = find_placeholder(slide, PP_PLACEHOLDER.BODY, hint_idx=2, hint_name="Caption")

                if title_ph and title_text:
                    align = config.get('alignment_layout_8_title', 'LEFT')
                    populate_text_frame(title_ph.text_frame, title_text, config['default_title_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config, alignment=align, vertical_anchor=default_vert_anchor)
                caption_text = slide_info.get('caption_text')
                if caption_ph and caption_text:
                    align = config.get('alignment_layout_8_caption', 'LEFT')
                    populate_text_frame(caption_ph.text_frame, caption_text, config['default_caption_font_size_pt'], config['default_font_name'], config['default_font_color_rgb'], config, alignment=align, vertical_anchor=default_vert_anchor) # Use caption size
                pic_desc = slide_info.get('picture_description')
                if picture_ph:
                    print(f"    Picture placeholder found (idx {picture_ph.placeholder_format.idx}).")
                    # Image insertion logic (requires image path in JSON)
                    image_path = slide_info.get("image_path")
                    if image_path and os.path.exists(image_path):
                        try:
                            # picture_ph.insert_picture(image_path) # Uncomment when image paths are provided
                            print(f"      (Simulated) Inserting image from {image_path}")
                        except Exception as img_e:
                            print(f"      Error inserting image '{image_path}': {img_e}")
                    elif image_path:
                         print(f"      Warning: Image file not found: {image_path}")
                    elif pic_desc:
                         # Optionally add description text if no image path provided and placeholder supports text
                         if picture_ph.has_text_frame:
                             populate_text_frame(picture_ph.text_frame, f"[{pic_desc}]", config['default_body_font_size_pt']-4, config['default_font_name'], "888888", config, alignment="CENTER", vertical_anchor="MIDDLE")
                         else:
                             print(f"      Picture description provided ('{pic_desc}') but no image path, and placeholder cannot hold text.")

            # --- Populate Notes ---
            if notes_text:
                try:
                    notes_slide = slide.notes_slide
                    text_frame = notes_slide.notes_text_frame
                    align = config.get('alignment_notes', 'LEFT')
                    populate_text_frame(text_frame, notes_text,
                                        config['default_notes_font_size_pt'],
                                        config['default_font_name'],
                                        config['default_font_color_rgb'], config,
                                        alignment=align, vertical_anchor=default_vert_anchor)
                except Exception as notes_err:
                    print(f"  Warning: Failed to add notes to slide {i+1}: {notes_err}")


        except (ValueError, TypeError, KeyError, AttributeError, IndexError) as data_err:
            print(f"DATA/LOGIC ERROR processing slide {i+1}. Data: {slide_info}. Error: {data_err}")
            traceback.print_exc() # Show details for debugging
        except Exception as e:
            print(f"UNEXPECTED FATAL ERROR processing slide {i+1}. Data: {slide_info}. Error: {e}")
            traceback.print_exc()


    # --- Save Presentation ---
    try:
        prs.save(output_filename)
        print(f"\nPresentation saved successfully as {output_filename}")
    except IOError as e:
        print(f"Error saving presentation to '{output_filename}': {e}. Check permissions or path.")
    except Exception as e:
        print(f"Unexpected error saving presentation: {e}")
        traceback.print_exc()


# --- Main Execution Function ---
def create_ppt_without_template(json_input, config_dict: Dict[str, Any], output_ppt_file='generated_presentation.pptx'):
    """
    Loads config, parses JSON data from a string, and creates a PowerPoint file.
    Expects json_input to be a string containing a JSON object like: {"slides": [...]}.
    """
    app_config = config_dict
    # Start with default configuration
    final_config = DEFAULT_CONFIG.copy()

    # Override defaults with provided config_dict values where keys match
    if isinstance(config_dict, dict):
        for key, value in config_dict.items():
            if key in final_config:
                final_config[key] = value
            else:
                print(f"Warning: Key '{key}' in config_dict not found in default configuration. Ignoring.")
    elif config_dict is not None:
         print(f"Warning: Invalid config_dict provided (type: {type(config_dict)}). Using default configuration.")
    

    # Use the merged configuration for this run
    app_config = final_config

    slide_list = None
    try:
        loaded_data = json.loads(json_input)
        # Expecting a dictionary with a 'slides' key containing a list
        if isinstance(loaded_data, dict) and 'slides' in loaded_data:
            slide_list = loaded_data['slides']
            if not isinstance(slide_list, list):
                 print(f"Error: 'slides' key found in JSON, but its value is not a list (type: {type(slide_list)}).")
                 slide_list = None # Prevent processing invalid data
        else:
            print("Error: JSON input must be a dictionary containing a 'slides' key with a list value.")
            # Example structure: '{ "slides": [ { "layout_idx": 0, "title": "..." }, ... ] }'
            slide_list = None

    except json.JSONDecodeError as e:
        print(f"Error decoding JSON input string: {e}")
        slide_list = None
    except Exception as e:
        print(f"An unexpected error occurred loading JSON from string: {e}")
        slide_list = None

    # Create Presentation only if slide_list was loaded correctly
    if slide_list is not None: # Check for None (empty list is valid)
        create_ppt(slide_list, output_ppt_file, app_config)
    else:
        print(f"Presentation creation skipped because slide data failed to load or parse correctly from the JSON input.")