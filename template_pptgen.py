import json
import os
import re
import traceback
from typing import Any, Dict, List, Optional, Tuple, Union

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.slide import SlideLayout
from pptx.util import Pt

from placeholders_utils import find_placeholder_id, get_placeholder_details
from image_handler import get_image # Added import

# Maps standard layout indices to descriptive names
LAYOUT_IDXs = {
    0: "TITLE_SLIDE",
    1: "TITLE_AND_CONTENT",
    2: "SECTION_HEADER",
    3: "TWO_CONTENT",
    4: "COMPARISON",
    5: "TITLE_ONLY",
    6: "BLANK",
    7: "CONTENT_WITH_CAPTION",
    8: "PICTURE_WITH_CAPTION",
}

# Default mapping of descriptive names to layout indices for PowerPoint template
# Can be overridden by 'layout_mapping' in the config file
DEFAULT_LAYOUT_MAPPING = {
    "TITLE_SLIDE": 0,
    "TITLE_AND_CONTENT": 1,
    "SECTION_HEADER": 2,
    "TWO_CONTENT": 3,
    "COMPARISON": 4,
    "TITLE_ONLY": 5,
    "BLANK": 6,
    "CONTENT_WITH_CAPTION": 7,
    "PICTURE_WITH_CAPTION": 8,
}

# --- Helper Functions ---

def _parse_color(rgb_string: Optional[str]) -> Optional[RGBColor]:
    """Converts a hex RGB string (e.g., "FF0000") to an RGBColor object."""
    if rgb_string and len(rgb_string) == 6:
        try:
            return RGBColor.from_string(rgb_string)
        except ValueError:
            print(f"Warning: Invalid RGB color string '{rgb_string}'. Using default.")
    return None

def _get_alignment_enum(align_str: Optional[str]) -> PP_ALIGN:
    """Converts an alignment string to a PP_ALIGN enum."""
    align_map = {
        "LEFT": PP_ALIGN.LEFT,
        "CENTER": PP_ALIGN.CENTER,
        "RIGHT": PP_ALIGN.RIGHT,
        "JUSTIFY": PP_ALIGN.JUSTIFY,
        "DISTRIBUTE": PP_ALIGN.DISTRIBUTE,
        "THAI_DISTRIBUTE": PP_ALIGN.THAI_DISTRIBUTE,
    }
    return align_map.get(str(align_str).upper(), PP_ALIGN.LEFT) # Default to LEFT

def _get_vertical_anchor_enum(anchor_str: Optional[str]) -> MSO_ANCHOR:
    """Converts a vertical anchor string to an MSO_ANCHOR enum."""
    anchor_map = {
        "TOP": MSO_ANCHOR.TOP,
        "MIDDLE": MSO_ANCHOR.MIDDLE,
        "BOTTOM": MSO_ANCHOR.BOTTOM,
    }
    return anchor_map.get(str(anchor_str).upper(), MSO_ANCHOR.TOP) # Default to TOP

def _get_value_from_keys(data_dict: Dict[str, Any], keys_list: List[str]) -> Optional[Any]:
    """
    Retrieves the value from a dictionary corresponding to the first key found
    from a list of potential keys.
    """
    for key in keys_list:
        if key in data_dict:
            return data_dict[key]
    return None

def _set_font_defaults(run, config: Dict[str, Any], size_pt: Optional[float] = None, force_color: bool = False, force_name: bool = False):
    """Applies font settings from config to a run, minimizing overrides."""
    font = run.font

    if force_name and config.get('default_font_name'):
        font.name = config['default_font_name']

    # Determine if font size override should be applied
    apply_size_override = False
    if size_pt is not None:
        is_using_template = config.get('using_custom_template', False)
        allow_override_with_template = config.get('apply_font_size_overrides_with_template', False)
        if not is_using_template or allow_override_with_template:
            apply_size_override = True

    if apply_size_override:
        font.size = Pt(size_pt)

    color_rgb = _parse_color(config.get('default_font_color_rgb'))
    if force_color and color_rgb:
        font.color.rgb = color_rgb

def _apply_markdown_to_run(run, text: str) -> str:
    """
    Applies markdown formatting (**bold**, *italic*, <u>underline</u>) to a run
    by stripping markers and setting font properties. Returns the stripped text.
    """
    run_bold = False
    run_italic = False
    run_underline = False
    current_text = text

    # Iteratively strip outer formatting markers and set flags
    processed = True
    while processed:
        processed = False
        # Check for bold
        if current_text.startswith('**') and current_text.endswith('**') and len(current_text) > 4:
            current_text = current_text[2:-2]
            run_bold = True
            processed = True
        # Check for italic
        elif current_text.startswith('*') and current_text.endswith('*') and len(current_text) > 2:
            current_text = current_text[1:-1]
            run_italic = True
            processed = True
        # Check for underline
        elif current_text.startswith('<u>') and current_text.endswith('</u>') and len(current_text) > 7:
            current_text = current_text[3:-4]
            run_underline = True
            processed = True

    run.font.bold = run_bold
    run.font.italic = run_italic
    run.font.underline = run_underline
    return current_text

def _add_formatted_text(text_frame, text: str, config: Dict[str, Any], default_size_pt: Optional[float] = None):
    """Adds text to a text frame, handling markdown combinations."""
    text_frame.clear()
    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    p.text = "" # Ensure paragraph is empty

    # Regex to split text by markdown markers, keeping the markers as delimiters
    parts = re.split(r'(\*\*.*?\*\*|\*.*?\*|<u>.*?</u>)', text)

    force_name = config.get('force_default_font_name', False)
    force_color = config.get('force_default_font_color', False)

    for part in filter(None, parts): # Process non-empty parts
        run = p.add_run()
        stripped_text = _apply_markdown_to_run(run, part)
        run.text = stripped_text
        _set_font_defaults(run, config, size_pt=default_size_pt, force_color=force_color, force_name=force_name)

    # Clean up potentially empty first paragraph if necessary
    if len(text_frame.paragraphs) > 1 and not text_frame.paragraphs[0].runs and not text_frame.paragraphs[0].text.strip():
         first_p = text_frame.paragraphs[0]._p # Access internal element
         text_frame._txBody.remove(first_p) # Remove the empty paragraph element

def _add_formatted_text_to_placeholder(slide, placeholder_idx: int, text: str, config: Dict[str, Any], default_size_pt: Optional[float], default_align: str = "LEFT", default_vert_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP) -> bool:
    """Helper to add formatted text to a specific placeholder index."""
    try:
        shape = slide.placeholders[placeholder_idx]
        if shape.has_text_frame:
            tf = shape.text_frame
            _add_formatted_text(tf, text, config, default_size_pt)
            if tf.paragraphs:
                tf.paragraphs[0].alignment = _get_alignment_enum(default_align)
            shape.vertical_anchor = default_vert_anchor
            return True
        else:
            print(f"Warning: Placeholder {placeholder_idx} exists but has no text frame.")
    except KeyError:
        print(f"Warning: Placeholder with index {placeholder_idx} not found in this slide's layout.")
    except Exception as e:
         print(f"Error processing placeholder {placeholder_idx} for formatted text: {e}")
    return False

def _add_bulleted_content(slide, placeholder_idx: int, content_list: Union[str, List[str]], config: Dict[str, Any], default_align: str = "LEFT", default_vert_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP) -> bool:
    """Adds bulleted/indented content, handling markdown combinations."""
    try:
        body_shape = slide.placeholders[placeholder_idx]
        if body_shape.has_text_frame:
            tf = body_shape.text_frame
            tf.clear()
            body_align = _get_alignment_enum(default_align)
            body_shape.vertical_anchor = default_vert_anchor

            if isinstance(content_list, str):
                content_list = [content_list]
            if not content_list: # Handle empty list
                return True # Nothing to add, but not an error

            # --- Font Size Calculation ---
            base_font_size_override = config.get('default_body_font_size_pt')
            smaller_font_size_override = config.get('smaller_content_font_size_pt')
            min_font_size = config.get('min_font_size_pt', 10)
            reduction_per_level = config.get('indent_level_font_size_reduction_pt', 2)
            use_smaller_font_override = False

            if config.get('enable_dynamic_body_font_size', False):
                item_count = len(content_list)
                char_count = sum(len(str(item)) for item in content_list)
                if item_count > config.get('dynamic_size_item_count_threshold', 6) or \
                   char_count > config.get('dynamic_size_char_count_threshold', 400):
                    use_smaller_font_override = True
                    print(f"  Applying smaller font size override for body content on slide.")

            current_base_size_override = smaller_font_size_override if use_smaller_font_override else base_font_size_override
            force_name = config.get('force_default_font_name', False)
            force_color = config.get('force_default_font_color', False)
            # --- End Font Size Calculation ---

            # Clear the default paragraph if it exists and is empty
            if tf.paragraphs and not tf.paragraphs[0].runs and not tf.paragraphs[0].text.strip():
                 first_p = tf.paragraphs[0]._p
                 tf._txBody.remove(first_p)

            for item in content_list:
                original_item_str = str(item) # Keep original for space counting
                item_str = original_item_str.strip() # Use stripped for content
                leading_spaces = len(original_item_str) - len(original_item_str.lstrip(' '))
                level = leading_spaces // 2 # Assuming 2 spaces per indent level

                p = tf.add_paragraph()
                p.alignment = body_align
                p.level = level

                # Calculate the specific size for this paragraph's runs
                run_font_size = None
                if current_base_size_override is not None:
                    run_font_size = max(min_font_size, current_base_size_override - (level * reduction_per_level))

                # Apply Markdown Formatting (with combinations)
                parts = re.split(r'(\*\*.*?\*\*|\*.*?\*|<u>.*?</u>)', item_str)
                for part in filter(None, parts):
                    run = p.add_run()
                    stripped_text = _apply_markdown_to_run(run, part)
                    run.text = stripped_text
                    _set_font_defaults(run, config, size_pt=run_font_size, force_color=force_color, force_name=force_name)

            tf.word_wrap = True
            return True
        else:
            print(f"Warning: Placeholder {placeholder_idx} exists but has no text frame for content.")
    except KeyError:
         print(f"Warning: Placeholder with index {placeholder_idx} not found in this slide's layout (for bulleted content).")
    except Exception as e:
         print(f"Error processing bulleted content for placeholder {placeholder_idx}: {e}")
         # traceback.print_exc() # Uncomment for debugging specific bullet errors
    return False

def _handle_image_search_and_insert(slide, placeholder_idx: int, image_query: str, output_path: str, config: Dict[str, Any], slide_num: int) -> bool:
    """
    Helper function to search for an image and insert it into a slide placeholder.
    
    Args:
        slide: The slide object to insert the image into
        placeholder_idx: The index of the placeholder to insert the image into
        image_query: The search query for the image
        output_path: The path where the presentation will be saved (used to determine image save location)
        config: The configuration dictionary
        slide_num: The slide number (for logging)
        
    Returns:
        bool: True if image was successfully inserted, False otherwise
    """
    if not image_query or placeholder_idx is None:
        return False

    try:
        # Create images directory next to the output presentation
        image_save_dir = os.path.join(os.path.dirname(output_path) or '.', 'images')
        print(f"  Slide {slide_num}: Attempting to fetch image using query: '{image_query}'")
        
        # Search for and download the image
        pil_img, fetched_img_path, aspect_ratio, img_url = get_image(image_query, save_dir=image_save_dir)
        
        if not (fetched_img_path and os.path.exists(fetched_img_path)):
            if pil_img:
                print(f"Warning: Slide {slide_num}: Fetched image but failed to save it to disk. Cannot insert.")
            else:
                print(f"Warning: Slide {slide_num}: Failed to fetch image using query '{image_query}'.")
            return False

        # Insert the image into the placeholder
        try:
            picture_placeholder = slide.placeholders[placeholder_idx]
            picture_placeholder.insert_picture(fetched_img_path)
            print(f"  Inserted image '{fetched_img_path}' into placeholder {placeholder_idx}.")
            return True
        except Exception as img_err:
            print(f"Warning: Slide {slide_num}: Failed to insert image '{fetched_img_path}' into placeholder {placeholder_idx}. Error: {img_err}.")
            # Add fallback text if image insertion fails
            _add_formatted_text_to_placeholder(
                slide, placeholder_idx,
                f"[Image: {image_query} - error]",
                config, config.get('default_body_font_size_pt'),
                "CENTER", MSO_ANCHOR.MIDDLE
            )
            return False

    except Exception as e:
        print(f"Error in image search and insertion for slide {slide_num}: {e}")
        return False

# --- Main Functions ---

def extract_layouts(template_path: Optional[str], config_dict: Dict[str, Any]) -> Tuple[Optional[Presentation], Dict[str, SlideLayout], Dict[str, Any]]:
    """
    Loads/creates a presentation, extracts slide layouts based on the provided config dictionary.
    """
    # Use the provided configuration dictionary directly
    config = config_dict.copy() # Work with a copy to avoid modifying the original dict

    try:
        # Update default layout mapping with custom mappings from config
        current_theme = config.get('current_theme', 'default')
        custom_mapping = config.get('templates', {}).get(current_theme, {}).get('layout_mapping', {})
        layout_mapping_to_use = DEFAULT_LAYOUT_MAPPING.copy()
        layout_mapping_to_use.update(custom_mapping)
        print(f"Using theme: {current_theme}")

    except Exception as e:
        raise Exception(f"Error processing configuration dictionary: {e}") from e

    # Load or create presentation
    prs = None
    config['using_custom_template'] = False # Track if a template was used
    if template_path and os.path.exists(template_path):
        try:
            prs = Presentation(template_path)
            config['using_custom_template'] = True
            print(f"Loaded presentation from template: {template_path}")

            # --- Remove existing slides from the loaded template ---
            if prs.slides and len(prs.slides._sldIdLst) > 0:
                print(f"Removing {len(prs.slides)} existing slide(s) from template...")
                slide_id_list = prs.slides._sldIdLst
                # Iterate backwards to avoid index issues when removing
                for i in range(len(slide_id_list) - 1, -1, -1):
                    slide_id = slide_id_list[i]
                    try:
                        rId = slide_id.rId
                        prs.part.drop_rel(rId)
                        del slide_id_list[i]
                    except Exception as remove_err:
                        rId_str = slide_id.rId if hasattr(slide_id, 'rId') else 'N/A'
                        print(f"  Warning: Could not remove slide with rId {rId_str}: {remove_err}")
                print("Finished removing existing slides.")
            # --- End of slide removal ---
        except Exception as e:
            raise Exception(f"Error loading presentation from {template_path}: {e}") from e
    else:
        if template_path:
             print(f"Warning: Template path specified but not found: {template_path}. Creating default presentation.")
        else:
             print("No template path specified. Creating default presentation.")
        prs = Presentation()

    # Extract layouts based on the effective mapping indices
    layouts_by_name: Dict[str, SlideLayout] = {}
    available_layouts = prs.slide_layouts
    num_available = len(available_layouts)
    print(f"Template contains {num_available} layouts.")

    for name, index in layout_mapping_to_use.items():
        if 0 <= index < num_available:
            layout = available_layouts[index]
            layouts_by_name[name] = layout
            # print(f"  Mapping '{name}' to layout at index {index} (Template Name: '{layout.name}')")
        else:
            print(f"  Warning: Layout index {index} for '{name}' is out of bounds (0-{num_available-1}). Skipping this layout mapping.")

    print(f"Mapped {len(layouts_by_name)} layouts based on configuration.")
    return prs, layouts_by_name, config


def create_ppt_with_template(json_string: Union[str, Dict], output_path: str, config_dict: Dict[str, Any], template_path: Optional[str] = None):
    """
    Creates a PowerPoint presentation based on JSON data, using a template and a configuration dictionary.
    """
    prs, layouts_by_name, config = None, {}, {}
    try:
        # Pass the config dictionary directly to extract_layouts
        prs, layouts_by_name, config = extract_layouts(template_path, config_dict)

        if not prs:
            raise RuntimeError("Failed to load or create a presentation object.")

        # --- Pre-load common config values ---
        default_vert_anchor = _get_vertical_anchor_enum(config.get("default_vertical_anchor", "TOP"))
        title_size_override = config.get('default_title_font_size_pt')
        subtitle_size_override = config.get('default_subtitle_font_size_pt')
        body_size_override = config.get('default_body_font_size_pt')
        notes_size_override = config.get('default_notes_font_size_pt')
        quote_body_size_override = config.get('default_quote_body_font_size_pt')
        # --- End Pre-load ---

        json_data = json.loads(json_string) if isinstance(json_string, str) else json_string
        if not isinstance(json_data, dict):
            raise ValueError("Invalid JSON data format. Expected a dictionary (object).")

        slides_data = json_data.get("slides", [])
        if not isinstance(slides_data, list):
             raise ValueError("Invalid JSON data format. 'slides' key must contain a list.")

        for i, slide_info in enumerate(slides_data):
            if not isinstance(slide_info, dict):
                print(f"Warning: Slide {i+1}: Invalid data format (expected dictionary). Skipping slide.")
                continue

            layout_idx = slide_info.get("layout_idx")
            layout_name = LAYOUT_IDXs.get(layout_idx) # Get standard name from index

            if layout_name is None or layout_name not in layouts_by_name:
                print(f"Warning: Slide {i+1}: Layout index {layout_idx} ('{layout_name}') not found or not mapped in template config. Skipping slide.")
                continue

            layout = layouts_by_name[layout_name]
            slide = prs.slides.add_slide(layout)
            # Get placeholder details for this specific layout
            # No need to modify this list during processing for this slide
            placeholders_info = get_placeholder_details(layout)
            print(f"Adding Slide {i+1}: Layout '{layout_name}' (Index {layout_idx}, Template: '{layout.name}')")

            # --- Helper function to populate a text placeholder ---
            def populate_text(key_in_json: str, placeholder_name: str, default_align: str, size_override: Optional[float], anchor: MSO_ANCHOR = default_vert_anchor):
                _idx = find_placeholder_id(placeholders_info, placeholder_name)
                if _idx is not None and key_in_json in slide_info:
                    align = config.get(f"alignment_layout_{layout_idx}_{placeholder_name}", default_align)
                    _add_formatted_text_to_placeholder(slide, _idx, slide_info[key_in_json], config, size_override, align, anchor)

            # --- Helper function to populate a content (bulleted) placeholder ---
            def populate_content(key_in_json: str, placeholder_name: str, default_align: str, anchor: MSO_ANCHOR = default_vert_anchor):
                 _idx = find_placeholder_id(placeholders_info, placeholder_name)
                 if _idx is not None and key_in_json in slide_info:
                     align = config.get(f"alignment_layout_{layout_idx}_{placeholder_name}", default_align)
                     _add_bulleted_content(slide, _idx, slide_info[key_in_json], config, align, anchor)

            # --- Populate Placeholders based on Layout Type ---
            # Layout 0: TITLE_SLIDE
            if layout_idx == 0:
                populate_text('title', 'title', "CENTER", title_size_override, MSO_ANCHOR.MIDDLE)
                populate_text('subtitle', 'subtitle', "CENTER", subtitle_size_override, MSO_ANCHOR.TOP) # Often anchored top

            # Layout 1: TITLE_AND_CONTENT
            elif layout_idx == 1:
                populate_text('title', 'title', "LEFT", title_size_override)
                populate_content('content', 'content', "LEFT")

            # Layout 2: SECTION_HEADER
            elif layout_idx == 2:
                 # Use more specific placeholder names if available in template
                populate_text('section_title', 'section_title', "LEFT", title_size_override)
                populate_text('section_description', 'section_description', "LEFT", body_size_override)
                 # Fallback to generic names if specific ones aren't found/used
                populate_text('title', 'title', "LEFT", title_size_override) # If 'section_title' wasn't populated
                populate_text('content', 'content', "LEFT", body_size_override) # If 'section_description' wasn't populated

            # Layout 3: TWO_CONTENT
            elif layout_idx == 3:
                populate_text('title', 'title', "LEFT", title_size_override)
                populate_content('left_content', 'left_content', "LEFT")
                populate_content('right_content', 'right_content', "LEFT")

            # Layout 4: COMPARISON
            elif layout_idx == 4:
                populate_text('title', 'title', "LEFT", title_size_override)

                # --- Left Side ---
                left_header_value = _get_value_from_keys(slide_info, ['left_header', 'left_title', 'left_heading'])
                left_content_value = _get_value_from_keys(slide_info, ['left_content', 'left_body', 'left_comparison_content'])
                left_header_idx = find_placeholder_id(placeholders_info, 'left_header')
                left_content_idx = find_placeholder_id(placeholders_info, 'left_content')
                left_header_processed = False

                # Try dedicated header placeholder first
                if left_header_value is not None and left_header_idx is not None:
                    align = config.get(f"alignment_layout_{layout_idx}_left_header", "CENTER")
                    if _add_formatted_text_to_placeholder(slide, left_header_idx , f"**{left_header_value.strip()}**", config, body_size_override, align, default_vert_anchor):
                        left_header_processed = True

                # Prepend header to content if not processed and content placeholder exists
                if not left_header_processed and left_header_value is not None and left_content_idx is not None:
                    header_text = f"**{left_header_value.strip()}**"
                    if left_content_value is None: left_content_value = [header_text]
                    elif isinstance(left_content_value, list): left_content_value.insert(0, header_text)
                    else: left_content_value = [header_text, str(left_content_value)]

                # Add content (possibly with header) to content placeholder
                if left_content_value is not None and left_content_idx is not None:
                     if left_content_idx != left_header_idx or not left_header_processed: # Avoid double-populating if same idx
                         align = config.get(f"alignment_layout_{layout_idx}_left_content", "LEFT")
                         _add_bulleted_content(slide, left_content_idx, left_content_value, config, align, default_vert_anchor)

                # --- Right Side (Similar logic) ---
                right_header_value = _get_value_from_keys(slide_info, ['right_header', 'right_title', 'right_heading'])
                right_content_value = _get_value_from_keys(slide_info, ['right_content', 'right_body', 'right_comparison_content'])
                right_header_idx = find_placeholder_id(placeholders_info, 'right_header')
                right_content_idx = find_placeholder_id(placeholders_info, 'right_content')
                right_header_processed = False

                if right_header_value is not None and right_header_idx is not None:
                    align = config.get(f"alignment_layout_{layout_idx}_right_header", "CENTER")
                    if _add_formatted_text_to_placeholder(slide, right_header_idx, f"**{right_header_value.strip()}**", config, body_size_override, align, default_vert_anchor):
                        right_header_processed = True

                if not right_header_processed and right_header_value is not None and right_content_idx is not None:
                    header_text = f"**{right_header_value.strip()}**"
                    if right_content_value is None: right_content_value = [header_text]
                    elif isinstance(right_content_value, list): right_content_value.insert(0, header_text)
                    else: right_content_value = [header_text, str(right_content_value)]

                if right_content_value is not None and right_content_idx is not None:
                     if right_content_idx != right_header_idx or not right_header_processed:
                         align = config.get(f"alignment_layout_{layout_idx}_right_content", "LEFT")
                         _add_bulleted_content(slide, right_content_idx, right_content_value, config, align, default_vert_anchor)


            # Layout 5: TITLE_ONLY (Often used for quotes or section breaks)
            elif layout_idx == 5:
                populate_text('title', 'title', "CENTER", title_size_override, MSO_ANCHOR.MIDDLE)
                # Allow optional 'subtitle' or 'body' text in the main body placeholder
                body_text = slide_info.get('subtitle') or slide_info.get('body')
                if body_text:
                    _idx = find_placeholder_id(placeholders_info, 'body') # Or a more specific name if template has one
                    if _idx is not None:
                         align = config.get(f"alignment_layout_{layout_idx}_body", "CENTER")
                         _add_formatted_text_to_placeholder(slide, _idx, body_text, config, quote_body_size_override, align, MSO_ANCHOR.MIDDLE)


            # Layout 6: BLANK (Allow adding content if specified, though unusual)
            elif layout_idx == 6:
                populate_text('title', 'title', "LEFT", title_size_override) # If a title placeholder exists
                populate_content('left_content', 'left_content', "LEFT") # If content placeholders exist
                populate_content('right_content', 'right_content', "LEFT")


            # Layout 7: CONTENT_WITH_CAPTION
            elif layout_idx == 7:
                populate_text('title', 'title', "LEFT", title_size_override)
                # caption_text is the main content area, potentially with bullets
                populate_content('caption_text', 'caption_text', "LEFT", default_vert_anchor)

                # object_description is used as a query for an image
                image_query = slide_info.get('object_description')
                placeholder_idx_img = find_placeholder_id(placeholders_info, 'object_description') # Common name
                if placeholder_idx_img is None:
                    placeholder_idx_img = find_placeholder_id(placeholders_info, 'picture_description') # Alternative name
                if placeholder_idx_img is None:
                    placeholder_idx_img = find_placeholder_id(placeholders_info, 'media') # Another alternative

                if image_query and placeholder_idx_img is not None:
                    success = _handle_image_search_and_insert(slide, placeholder_idx_img, image_query, output_path, config, i+1)
                    if not success and 'image_description' in slide_info:
                        # Fallback to showing the description if image search/insertion failed
                        populate_text('image_description', 'picture', "CENTER", body_size_override, MSO_ANCHOR.MIDDLE)
                elif image_query and placeholder_idx_img is None:
                    print(f"Warning: Slide {i+1}: Image query '{image_query}' provided for layout 7, but no suitable picture/object placeholder found.")


            # Layout 8: PICTURE_WITH_CAPTION
            elif layout_idx == 8:
                populate_content('caption_text', 'caption_text', "LEFT") # Text content

                # --- Handle Picture ---
                placeholder_idx_img = find_placeholder_id(placeholders_info, 'picture') # Common name
                if placeholder_idx_img is None:
                    placeholder_idx_img = find_placeholder_id(placeholders_info, 'picture_description') # Alternative name

                # First try image_path if provided
                img_path_to_use = slide_info.get('image_path')
                if img_path_to_use and os.path.exists(img_path_to_use):
                    try:
                        picture_placeholder = slide.placeholders[placeholder_idx_img]
                        picture_placeholder.insert_picture(img_path_to_use)
                        print(f"  Inserted image '{img_path_to_use}' into placeholder {placeholder_idx_img}.")
                    except Exception as img_err:
                        print(f"Warning: Slide {i+1}: Failed to insert image '{img_path_to_use}' into placeholder {placeholder_idx_img}. Error: {img_err}.")
                        img_path_to_use = None  # Clear this so we try image search next
                
                # If no image_path or it failed, try picture_description as search query
                if not img_path_to_use:
                    image_query = slide_info.get('picture_description')
                    if image_query and placeholder_idx_img is not None:
                        success = _handle_image_search_and_insert(slide, placeholder_idx_img, image_query, output_path, config, i+1)
                        if not success and 'image_description' in slide_info:
                            # Fallback to showing the description if image search/insertion failed
                            populate_text('image_description', 'picture', "CENTER", body_size_override, MSO_ANCHOR.MIDDLE)
                    elif image_query:
                        print(f"Warning: Slide {i+1}: Image query provided but no suitable picture placeholder found in layout {layout_idx}.")


            # --- Handle Notes ---
            if 'notes' in slide_info and slide_info['notes']: # Check if notes exist and are not empty
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    tf = notes_slide.notes_text_frame
                    if tf:
                        _add_formatted_text(tf, slide_info['notes'], config, notes_size_override)
                        if tf.paragraphs:
                            tf.paragraphs[0].alignment = _get_alignment_enum(config.get("alignment_notes", "LEFT"))
                    else:
                        print(f"Warning: Slide {i+1}: Notes provided but notes text frame not found in notes slide.")
                else:
                     print(f"Warning: Slide {i+1}: Notes provided but slide does not have a notes slide.")


        # --- Save Presentation ---
        print(f"\nSaving presentation to {output_path}")
        prs.save(output_path)
        print("Presentation saved successfully.")

    except FileNotFoundError as e:
        print(f"Error: Required file not found. {e}")
    except (ValueError, TypeError) as e:
         print(f"Error: Invalid data format provided. {e}")
    except KeyError as e:
         print(f"Error: Missing expected key in data or configuration: {e}")
    except ImportError:
        print("Error: Could not import required libraries. Make sure 'python-pptx', 'PyYAML', and 'pptxexp' are installed correctly.")
        print("       Try: pip install python-pptx PyYAML pptxexp")
    except Exception as e:
        print(f"An unexpected error occurred: {e}")
        traceback.print_exc() # Print detailed traceback for debugging
